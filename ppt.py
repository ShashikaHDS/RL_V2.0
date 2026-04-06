from pptx import Presentation

# Create a new PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Rectangle Reduction Environment for Multi-Robot RL"
subtitle.text = "Using Stable Baselines3 and Custom Pygame Visualization"

# Slide 2: Problem Statement
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Problem Statement"
content.text = (
    "Objective: Train multiple robots to minimize the area of the rectangle "
    "formed by their positions.\n\n"
    "Constraints:\n"
    "- Robots must explore a grid-based map.\n"
    "- Avoid obstacles while optimizing the rectangle area.\n"
    "- Observe the environment using a limited field of view."
)

# Slide 3: Approach
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Approach"
content.text = (
    "Grid-Based Environment:\n"
    "- The map is represented as a 2D grid.\n"
    "- Each cell can be free space, an obstacle, or occupied by a robot.\n\n"
    "Actions:\n"
    "- Move Up, Down, Left, Right, or Stay.\n\n"
    "Reward Function:\n"
    "- +10: Achieving the optimal rectangle area (3x3).\n"
    "- +1: Reducing the rectangle area.\n"
    "- -1: Increasing the rectangle area.\n"
    "- -2: Collision with an obstacle.\n\n"
    "Tools:\n"
    "- Stable Baselines3 for reinforcement learning.\n"
    "- Pygame for real-time visualization."
)

# Slide 4: Environment Design
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Environment Design"
content.text = (
    "Environment Variables:\n"
    "- Grid Map: Defines the environment layout.\n"
    "- Number of Robots: Configurable parameter.\n"
    "- Robot Positions: Randomly initialized within free cells.\n"
    "- Field of View: Determines the exploration area around robots.\n\n"
    "Custom Gym Environment:\n"
    "- Action Space: MultiDiscrete (5 actions per robot).\n"
    "- Observation Space: Known map and robot positions."
)

# Slide 5: Code Overview: Initialization
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Code Overview: Initialization"
content.text = (
    "class RectangleReductionEnv(gym.Env):\n"
    "    def __init__(self, grid_map, num_robots, robot_positions, field_of_view):\n"
    "        self.grid_map = np.array(grid_map)\n"
    "        self.num_robots = num_robots\n"
    "        self.robot_positions = np.array(robot_positions)\n"
    "        self.field_of_view = field_of_view\n"
    "        self.grid_size = self.grid_map.shape\n"
    "        self.known_map = np.zeros_like(self.grid_map)\n"
    "        self.action_space = gym.spaces.MultiDiscrete([5] * self.num_robots)\n"
    "        self.observation_space = gym.spaces.Dict({\n"
    "            'known_map': gym.spaces.Box(0, 1, shape=self.grid_map.shape, dtype=np.uint8),\n"
    "            'robot_positions': gym.spaces.Box(0, max(self.grid_size), shape=(self.num_robots, 2), dtype=np.int32),\n"
    "        })"
)

# Save the presentation


import os

save_dir = r"D:\MNEED\Thesis\ppt\x"
os.makedirs(save_dir, exist_ok=True)
presentation_path = os.path.join(save_dir, "rectangle_reduction_presentation.pptx")
presentation.save(presentation_path)

